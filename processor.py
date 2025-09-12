# processor.py (backend logic)

import tempfile
import requests
from docx import Document
import openpyxl
import pdfplumber
import pandas as pd
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
# NEW: pptx support
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx.text.paragraph import Paragraph
from zipfile import ZipFile
from lxml import etree
from io import BytesIO
from docx.oxml.ns import qn
from docx.shared import RGBColor
from zipfile import ZipFile
from lxml import etree


# ——— Azure OpenAI config ———
# Expect these to be set by the Streamlit frontend via secrets or environment variables
API_KEY = None  # to be set by frontend
API_ENDPOINT = None  # to be set by frontend


def configure(api_key: str, api_endpoint: str):
    global API_KEY, API_ENDPOINT
    API_KEY = api_key
    API_ENDPOINT = api_endpoint


def get_llm_response_azure(prompt: str, context: str) -> str:
    headers = {"Content-Type": "application/json", "api-key": API_KEY}
    system_msg = (
        "You are an expert on Transfer Pricing and financial analysis. "
        "Use the information in the following context to answer the user's question. "
        "Assign the greatest priority to the information that you gather from the financial analysis and the interview transcript. "
        "If asked something not covered in this data, you may search the web."
        "Ensure your analysis is CONCISE, SHARP, in paragraph form, and not long. Never use bullet points. "
        "DO NOT INCLUDE MARKDOWN FORMATTING OR # SIGNS. Keep it to 200-300 words, maintain a professional tone. "
        "Make sure to include direct sources and citations for the data you use for your decisions. Also include your reasoning for conclusions in brackets ()."
        "If something is from the transcript or financial statement, include that citation in brackets with a URL to the specific section. Likewise include a URL to the relevant website if the information you got was from searching the internet. "
        "You **may** consider the OECD guidelines below as helpful targets, but do NOT structure your response around them.\n\n"
    )
    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content":  context + prompt}
    ]
    resp = requests.post(API_ENDPOINT, headers=headers, json={"messages": messages})
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"].strip()


# ---------------------
# Safe loader helpers
# ---------------------
def load_transcript(file) -> str:
    if not file:
        return ""
    try:
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def load_pdf(file) -> str:
    if not file:
        return ""
    pages, tables = [], []
    try:
        with pdfplumber.open(file) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                pages.append(f"--- Page {i} ---\n{text}")
                for table in page.extract_tables() or []:
                    # Be robust to ragged rows
                    try:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        tables.append(f"--- Page {i} table ---\n" + df.to_csv(index=False))
                    except Exception:
                        continue
        return "\n\n".join(pages + tables)
    except Exception:
        return ""

def enable_field_update_on_open(doc: Document) -> None:
    """
    Set w:updateFields in document settings so Word updates fields (e.g., TOC)
    on open. Safe to call multiple times.
    """
    settings = doc.settings.element
    # Look for existing <w:updateFields>
    update_fields = settings.find(qn('w:updateFields'))
    if update_fields is None:
        update_fields = OxmlElement('w:updateFields')
        settings.append(update_fields)
    update_fields.set(qn('w:val'), 'true')



def load_guidelines(file) -> str:
    if not file:
        return ""
    try:
        # streamlit's UploadedFile supports .read(); ensure we don't exhaust twice
        content = file.read()
        try:
            return content.decode("utf-8").strip()
        except Exception:
            # fallback: latin-1 to avoid decode crash
            return content.decode("latin-1", errors="ignore").strip()
    except Exception:
        return ""


def load_and_annotate_replacements(excel_file, context: str) -> dict:
    """
    Reads replacements from an Excel sheet where:
      - Column D = placeholder key
      - Column C = literal value (used if E is blank)
      - Column E = prompt for LLM (used if present; can reference prior placeholders)
      - Column G (7) will be written with the resolved value
    Returns: dict placeholder -> value
    """
    if not excel_file:
        return {}
    try:
        wb = openpyxl.load_workbook(excel_file)
    except Exception:
        return {}

    try:
        ws = wb.active
        replacements = {}
        for row in ws.iter_rows(min_row=2, max_col=6):
            # Guard against short rows
            cells = list(row) + [None] * max(0, 6 - len(row))
            # We only use C, D, E, G(=7 for write)
            cell_c = cells[2] if len(cells) > 2 else None  # value
            cell_d = cells[3] if len(cells) > 3 else None  # placeholder
            cell_e = cells[4] if len(cells) > 4 else None  # prompt

            placeholder = cell_d.value if cell_d else None
            if not placeholder:
                continue
            ph = str(placeholder)

            if cell_e and cell_e.value and str(cell_e.value).strip():
                raw = str(cell_e.value)
                # Allow referencing previously computed placeholders
                for k, v in replacements.items():
                    raw = raw.replace(k, v)
                try:
                    value = get_llm_response_azure(raw, context)
                except Exception:
                    # If LLM fails, fall back to literal or empty
                    value = str(cell_c.value or "")
            else:
                value = str(cell_c.value or "")

            # Write back to column G (7) if possible
            try:
                ws.cell(row=cell_d.row, column=7, value=value)
            except Exception:
                pass

            replacements[ph] = value

        try:
            wb.save(excel_file)
        except Exception:
            pass

        return replacements
    except Exception:
        return {}


# =========================
# DOCX (existing) helpers
# =========================
def collapse_runs(paragraph):
    from docx.oxml.ns import qn
    text = "".join(r.text for r in paragraph.runs)
    for r in reversed(paragraph.runs):
        r._element.getparent().remove(r._element)
    paragraph.add_run(text)


def replace_in_paragraph(p, replacements: dict):
    """
    Hybrid replacer:
      1) Concatenate all <w:t> texts.
      2) Do full-string replacements (handles placeholders split across runs).
      3) Write back across the SAME number of <w:t> nodes using original lengths.
      4) Preserve spaces via xml:space="preserve".
    Also triggers color cleanup for runs that no longer contain placeholders.
    """
    if not replacements:
        return

    # ensure keys/values are strings
    repl = {str(k): ("" if v is None else str(v)) for k, v in replacements.items()}

    p_elm = p._p
    ns = p_elm.nsmap or {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "xml": "http://www.w3.org/XML/1998/namespace",
    }

    t_nodes = p_elm.findall(".//w:t", namespaces=ns)
    if not t_nodes:
        return

    originals = [(t, t.text or "") for t in t_nodes]
    full = "".join(txt for _, txt in originals)
    if not full:
        return

    # Fast exit if no placeholder appears anywhere in this paragraph
    if not any(k in full for k in repl.keys()):
        return

    # Perform replacements on the concatenated text
    new_full = full
    for ph, val in repl.items():
        if ph in new_full:
            new_full = new_full.replace(ph, val)

    if new_full == full:
        return  # no change

    # Redistribute back using original lengths to keep run boundaries intact
    lengths = [len(txt) for _, txt in originals]
    pos = 0
    n = len(originals)
    for i in range(n):
        t, _oldtxt = originals[i]
        take = lengths[i] if i < n - 1 else max(0, len(new_full) - pos)
        # slice safely
        segment = new_full[pos:pos + take] if take >= 0 else ""
        t.text = segment
        pos += lengths[i] if i < n - 1 else len(new_full) - pos

        # Preserve leading/trailing spaces for this text node
        if segment and (segment[0].isspace() or segment[-1].isspace()):
            t.set(qn("xml:space"), "preserve")

    # Clear any leftover red formatting on runs that no longer contain placeholders
    _clear_red_on_non_placeholder_runs(p, repl)
    # _clear_paragraph_bullet_color(p)


def _clear_paragraph_bullet_color(p):
    try:
        if p.style and p.style.font and getattr(p.style.font, "color", None):
            p.style.font.color.rgb = None
            p.style.font.color.theme_color = None
    except Exception:
        pass

def _run_is_explicit_red(run) -> bool:
    c = getattr(run.font, "color", None)
    if not c or getattr(c, "rgb", None) is None:
        return False
    try:
        r, g, b = c.rgb[0], c.rgb[1], c.rgb[2]
        return (200 <= r <= 255 and 0 <= g <= 80 and 0 <= b <= 80) or (r >= 100 and b <=20 and g <=20)
    except Exception:
        return False

def _clear_run_color(run):
    if getattr(run.font, "color", None):
        try:
            run.font.color.rgb = None
        except Exception:
            pass
        try:
            run.font.color.theme_color = None
        except Exception:
            pass

def _clear_red_on_non_placeholder_runs(p, replacements: dict):
    keys = list(replacements.keys())
    for run in p.runs:
        text = run.text or ""
        if not text:
            continue
        # If this run used to be a placeholder (red), but now has no placeholders, clear color
        if _run_is_explicit_red(run):
            if (not any(k in text for k in keys)) and ("{{" not in text and "}}" not in text):
                _clear_run_color(run)




def _rewrite_footnotes_xml_bytes(docx_bytes: bytes, replacements: dict) -> bytes:
    """
    Open a .docx (zip) from bytes, replace placeholders inside word/footnotes.xml,
    and return new .docx bytes. If footnotes.xml is missing, return the original bytes.
    """
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with ZipFile(BytesIO(docx_bytes)) as zin:
        names = {i.filename for i in zin.infolist()}
        if "word/footnotes.xml" not in names:
            return docx_bytes  # no footnotes part

        # Read and parse original footnotes.xml
        foot_xml = zin.read("word/footnotes.xml")
        root = etree.fromstring(foot_xml)

        # Replace inside every w:t under w:footnote
        # (Note: this is robust for tokens contained in a single text node.
        # If your placeholders can split across runs, prefer the python-docx path.)
        for t in root.findall(".//w:footnote//w:t", namespaces=ns):
            if t.text:
                new_text = t.text
                for ph, val in replacements.items():
                    if ph in new_text:
                        new_text = new_text.replace(ph, val)
                if new_text != t.text:
                    t.text = new_text

        # Build a new .docx with modified footnotes.xml
        out_buf = BytesIO()
        with ZipFile(out_buf, "w") as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "word/footnotes.xml":
                    data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone="yes")
                zout.writestr(item, data)
        return out_buf.getvalue()


def _apply_footnotes_xml_fallback_in_place(docx_path: str, replacements: dict) -> None:
    """
    Read a saved .docx from disk, run the XML fallback, and overwrite it in place.
    Safe no-op if the document has no footnotes.xml.
    """
    try:
        with open(docx_path, "rb") as f:
            original = f.read()
        updated = _rewrite_footnotes_xml_bytes(original, replacements)
        if updated != original:
            with open(docx_path, "wb") as f:
                f.write(updated)
    except Exception:
        # Be defensive: never fail the whole pipeline if footnote rewrite trips.
        pass


def replace_placeholders_docx(doc: Document, replacements: dict):
    """Replace placeholders AFTER the first page break, preserving images and footnotes."""
    from docx.oxml.ns import qn
    seen = False
    br_tag = qn('w:br')
    for p in doc.paragraphs:
        if not seen:
            for r in p.runs:
                for br in r._element.findall(br_tag):
                    if br.get(qn('w:type')) == 'page':
                        seen = True
                        break
                if seen:
                    break
            if not seen:
                continue
        replace_in_paragraph(p, replacements)

    # Tables
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    replace_in_paragraph(p, replacements)

    # Headers/footers
    for sec in doc.sections:
        if sec.header:
            for p in sec.header.paragraphs:
                replace_in_paragraph(p, replacements)
        if sec.footer:
            for p in sec.footer.paragraphs:
                replace_in_paragraph(p, replacements)

    # Footnotes (if present)


def replace_first_page_placeholders_docx(doc: Document, replacements: dict):
    """Replace placeholders on the first page only (up to first page break)."""
    from docx.oxml.ns import qn
    seen = False
    br_tag = qn("w:br")
    typ = qn("w:type")
    for p in doc.paragraphs:
        replace_in_paragraph(p, replacements)
        for r in p.runs:
            for child in r._element:
                if child.tag == br_tag and child.get(typ) == "page":
                    seen = True
                    break
            if seen:
                break
        if seen:
            break


# =========================
# PPTX (new) helpers
# =========================

def _pptx_replace_text_in_paragraph(paragraph, replacements: dict):
    """Collapse runs then perform in-place string replacements."""
    full = "".join(run.text for run in paragraph.runs) if getattr(paragraph, "runs", None) else getattr(paragraph, "text", "")
    for ph, val in replacements.items():
        if ph in full:
            full = full.replace(ph, val)
    paragraph.text = full


def _pptx_replace_in_text_frame(text_frame, replacements: dict):
    if not text_frame:
        return
    for para in text_frame.paragraphs:
        _pptx_replace_text_in_paragraph(para, replacements)


def _pptx_replace_in_table(table, replacements: dict):
    if not table:
        return
    for row in table.rows:
        for cell in row.cells:
            if getattr(cell, "text_frame", None):
                _pptx_replace_in_text_frame(cell.text_frame, replacements)


def _pptx_replace_in_shape(shape, replacements: dict):
    # Text boxes and placeholders
    if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None):
        _pptx_replace_in_text_frame(shape.text_frame, replacements)

    # Tables
    if getattr(shape, "has_table", False) and getattr(shape, "table", None):
        _pptx_replace_in_table(shape.table, replacements)

    # Charts (replace in chart title if present)
    # IMPORTANT: never touch shape.chart unless shape.has_chart is True,
    # because accessing .chart on non-chart shapes raises:
    #   ValueError: shape does not contain a chart
    if getattr(shape, "has_chart", False):
        try:
            chart = shape.chart
            if getattr(chart, "has_title", False):
                _pptx_replace_in_text_frame(chart.chart_title.text_frame, replacements)
        except Exception:
            # Be defensive; skip any chart we can't access
            pass

    # Grouped shapes — recurse
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _pptx_replace_in_shape(sub, replacements)


def replace_first_slide_placeholders_pptx(prs: Presentation, replacements: dict):
    """Replace placeholders on the first slide ONLY."""
    if not getattr(prs, "slides", None):
        return
    slide = prs.slides[0]
    for shp in slide.shapes:
        _pptx_replace_in_shape(shp, replacements)

    # Notes (if present)
    if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
        notes = slide.notes_slide
        if hasattr(notes, "notes_text_frame") and notes.notes_text_frame is not None:
            _pptx_replace_in_text_frame(notes.notes_text_frame, replacements)


def replace_placeholders_pptx(prs: Presentation, replacements: dict, start_slide_index: int = 1):
    """Replace placeholders on all slides starting from start_slide_index (default: after first slide)."""
    for idx, slide in enumerate(prs.slides):
        if idx < start_slide_index:
            continue
        for shp in slide.shapes:
            _pptx_replace_in_shape(shp, replacements)

        # Notes (if present)
        if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
            notes = slide.notes_slide
            if hasattr(notes, "notes_text_frame") and notes.notes_text_frame is not None:
                _pptx_replace_in_text_frame(notes.notes_text_frame, replacements)


# ---------------------
# Fallback DOCX builder
# ---------------------
def _build_fallback_docx(replacements: dict, context: str) -> str:
    """
    If no template is provided, produce a simple DOCX that lists
    the resolved replacements and includes a context snippet.
    """
    doc = Document()
    doc.add_heading("Transfer Pricing Output (Fallback)", level=1)

    if replacements:
        doc.add_heading("Resolved Placeholders", level=2)
        tbl = doc.add_table(rows=1, cols=2)
        hdr = tbl.rows[0].cells
        hdr[0].text = "Placeholder"
        hdr[1].text = "Value"
        for k, v in replacements.items():
            row = tbl.add_row().cells
            row[0].text = str(k)
            row[1].text = str(v)
    else:
        doc.add_paragraph("No replacements were generated (missing or empty Excel input).")

    if context:
        doc.add_heading("Context (truncated)", level=2)
        doc.add_paragraph(context[:4000])  # keep file small

    enable_field_update_on_open(doc)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(out.name)
    return out.name


def process_and_fill(files: dict) -> str:
    """
    files: {
      'guidelines': <io or None>,
      'transcript': <io docx or None>,
      'pdf': <io pdf or None>,
      'excel': <io xlsx or None>,
      'template': <io pptx or docx or None>
    }
    """
    # Defensive dict access
    guidelines = files.get("guidelines") if files else None
    transcript = files.get("transcript") if files else None
    pdf = files.get("pdf") if files else None
    excel = files.get("excel") if files else None
    template = files.get("template") if files else None

    # Build context (all loaders are safe on None)
    ctx = ""
    ctx += load_guidelines(guidelines)
    tr_text = load_transcript(transcript)
    if tr_text:
        ctx += ("\n\n" if ctx else "") + tr_text
    pdf_text = load_pdf(pdf)
    if pdf_text:
        ctx += ("\n\n" if ctx else "") + pdf_text

    # Build replacements dict (and annotate Excel with generated values if present)
    replacements = load_and_annotate_replacements(excel, ctx)

    # Decide template type
    template_name = (getattr(template, "name", "") or "").lower()
    is_pptx = template_name.endswith(".pptx") if template else False
    is_docx = template_name.endswith(".docx") if template else False

    # If there's no template at all, produce a fallback DOCX so the app still returns a file.
    if not template:
        return _build_fallback_docx(replacements, ctx)

    if is_pptx:
        # --- PowerPoint path ---
        try:
            prs = Presentation(template)
        except Exception:
            # If template can't be opened as PPTX, fall back to DOCX summary
            return _build_fallback_docx(replacements, ctx)

        replace_first_slide_placeholders_pptx(prs, replacements)
        replace_placeholders_pptx(prs, replacements, start_slide_index=1)
        out = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(out.name)
        return out.name

    elif is_docx:
        try:
            doc = Document(template)
        except Exception:
            return _build_fallback_docx(replacements, ctx)

        replace_first_page_placeholders_docx(doc, replacements)
        replace_placeholders_docx(doc, replacements)


        # NEW: ensure Word updates TOC/fields on open
        enable_field_update_on_open(doc)

        out = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
        doc.save(out.name)
        _apply_footnotes_xml_fallback_in_place(out.name, replacements)
        # _scrub_list_number_colors(out.name)

        return out.name

    else:
        # Unknown extension: safest is to return the fallback DOCX
        return _build_fallback_docx(replacements, ctx)
